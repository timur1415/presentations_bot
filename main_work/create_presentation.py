import json
import random
import textwrap

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def hex_to_rgb(hex_color):
    hex_color = hex_color.replace("#", "")

    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def get_circle_phrase(title, bullets):

    words = title.split()

    if len(words) >= 2:
        return f"{words[0]}\n{words[1]}"

    if bullets:

        bullet_words = bullets[0].split()

        if len(bullet_words) >= 2:
            return f"{bullet_words[0]}\n{bullet_words[1]}"

        return bullets[0][:20]

    return title


def add_circle_with_text(slide, style, text):

    size = random.uniform(1.5, 2.1)

    x = random.choice(
        [
            -0.5,
            11.0,
        ]
    )

    y = random.choice(
        [
            -0.3,
            5.2,
            2.8,
        ]
    )

    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x),
        Inches(y),
        Inches(size),
        Inches(size),
    )

    circle.fill.solid()
    circle.fill.fore_color.rgb = style["accent"]
    circle.fill.transparency = 0.12
    circle.line.fill.background()

    tf = circle.text_frame
    tf.clear()

    p = tf.paragraphs[0]

    p.text = text

    p.alignment = PP_ALIGN.CENTER

    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = style["text"]


def add_decor(slide, style, section=None):

    accent = style["accent"]

    count = random.choice([1, 1, 2])

    for _ in range(count):

        size = random.uniform(2.0, 4.2)

        x = random.choice([
            -1.2,
            11.0,
        ])

        y = random.choice([
            -1.0,
            5.3,
            random.uniform(0.8, 4.2),
        ])

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x),
            Inches(y),
            Inches(size),
            Inches(size),
        )

        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.fill.transparency = 0.82
        circle.line.fill.background()

    if random.random() > 0.7:

        diamond = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND,
            Inches(random.choice([0.8, 11.4])),
            Inches(random.choice([0.8, 5.8])),
            Inches(0.7),
            Inches(0.7),
        )

        diamond.fill.solid()
        diamond.fill.fore_color.rgb = accent
        diamond.fill.transparency = 0.84
        diamond.line.fill.background() 


def add_title(slide, title, style):

    box = slide.shapes.add_textbox(
        Inches(0.9),
        Inches(0.55),
        Inches(10.8),
        Inches(0.8),
    )

    p = box.text_frame.paragraphs[0]

    p.text = title

    p.font.size = Pt(26)

    p.font.bold = True

    p.font.color.rgb = style["title"]

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.9),
        Inches(1.18),
        Inches(random.uniform(1.6, 3.0)),
        Inches(0.03),
    )

    line.fill.solid()

    line.fill.fore_color.rgb = style["accent"]

    line.line.fill.background()


def add_bullets(slide, bullets, style):

    box = slide.shapes.add_textbox(
        Inches(1),
        Inches(1.75),
        Inches(10.7),
        Inches(4.8),
    )

    tf = box.text_frame

    tf.word_wrap = True

    for i, bullet in enumerate(bullets):

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        p.text = bullet

        p.font.size = Pt(18)

        p.font.color.rgb = style["text"]

        p.space_after = Pt(12)


def add_columns(slide, bullets, style):

    mid = len(bullets) // 2 + len(bullets) % 2

    left = bullets[:mid]

    right = bullets[mid:]

    for idx, col in enumerate([left, right]):

        x = 1.0 if idx == 0 else 6.3

        box = slide.shapes.add_textbox(
            Inches(x),
            Inches(1.75),
            Inches(4.8),
            Inches(4.8),
        )

        tf = box.text_frame

        tf.word_wrap = True

        for i, bullet in enumerate(col):

            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            p.text = bullet

            p.font.size = Pt(18)

            p.font.color.rgb = style["text"]

            p.space_after = Pt(10)


def add_paragraph(slide, bullets, style):

    text = " ".join(bullets)

    box = slide.shapes.add_textbox(
        Inches(1),
        Inches(1.75),
        Inches(10.8),
        Inches(4.8),
    )

    tf = box.text_frame

    p = tf.paragraphs[0]

    p.text = textwrap.fill(
        text,
        width=95
    )

    p.font.size = Pt(18)

    p.font.color.rgb = style["text"]


def add_content(slide, bullets, style):

    if not bullets:
        return

    variants = [
        "bullets",
        "columns",
        "paragraph",
        "mixed",
        "highlight",
    ]

    if len(bullets) <= 3:
        variant = random.choice(["bullets", "highlight"])

    elif len(bullets) >= 8:
        variant = random.choice(["columns", "paragraph"])

    else:
        variant = random.choice(variants)

    if variant == "bullets":
        add_bullets(slide, bullets, style)

    elif variant == "columns":
        add_columns(slide, bullets, style)

    elif variant == "paragraph":
        add_paragraph(slide, bullets, style)

    elif variant == "mixed":
        add_mixed(slide, bullets, style)

    else:
        add_highlight(slide, bullets, style)


def add_mixed(slide, bullets, style):

    if not bullets:
        return

    box = slide.shapes.add_textbox(
        Inches(1),
        Inches(1.8),
        Inches(10),
        Inches(0.8)
    )

    p = box.text_frame.paragraphs[0]

    p.text = bullets[0]

    p.font.size = Pt(24)

    p.font.bold = True

    p.font.color.rgb = style["text"]

    if len(bullets) > 1:

        box2 = slide.shapes.add_textbox(
            Inches(1),
            Inches(3),
            Inches(10),
            Inches(3),
        )

        tf = box2.text_frame

        for i, bullet in enumerate(bullets[1:]):

            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            p.text = f"• {bullet}"

            p.font.size = Pt(18)

            p.font.color.rgb = style["text"]


def add_highlight(slide, bullets, style):

    if not bullets:
        return

    box = slide.shapes.add_textbox(
        Inches(1.2),
        Inches(2.2),
        Inches(10),
        Inches(1.4)
    )

    p = box.text_frame.paragraphs[0]

    p.text = bullets[0]

    p.font.size = Pt(28)

    p.font.bold = True

    p.alignment = PP_ALIGN.CENTER

    p.font.color.rgb = style["text"]

    if len(bullets) > 1:

        box2 = slide.shapes.add_textbox(
            Inches(1.2),
            Inches(4.0),
            Inches(10),
            Inches(2)
        )

        tf = box2.text_frame

        for i, bullet in enumerate(bullets[1:]):

            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            p.text = bullet

            p.font.size = Pt(16)

            p.alignment = PP_ALIGN.CENTER

            p.font.color.rgb = style["text"]



def create_presentation(
    topic,
    slides_text,
    style_text=""
):

    prs = Presentation()

    data = json.loads(slides_text)

    style = {
        "background": hex_to_rgb(
            data["background_color"]
        ),
        "title": hex_to_rgb(
            data["title_color"]
        ),
        "text": hex_to_rgb(
            data["text_color"]
        ),
        "accent": hex_to_rgb(
            data["accent_color"]
        ),
    }

    for slide_data in data["slides"]:

        slide = prs.slides.add_slide(
            prs.slide_layouts[6]
        )

        fill = slide.background.fill

        fill.solid()

        fill.fore_color.rgb = style["background"]

        phrase = get_circle_phrase(
            slide_data["title"],
            slide_data["bullets"]
        )

        add_decor(
            slide,
            style,
            phrase
        )

        add_title(
            slide,
            slide_data["title"],
            style
        )

        add_content(
            slide,
            slide_data["bullets"],
            style
        )

    file_name = f"{topic}.pptx"

    prs.save(file_name)

    return file_name